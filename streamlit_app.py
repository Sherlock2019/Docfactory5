import streamlit as st
from docx import Document
from pptx import Presentation
from docx.shared import Inches
import pandas as pd
import re
from io import BytesIO
from tempfile import NamedTemporaryFile
from datetime import date
import os

st.set_page_config(page_title="🧩 Smart Docx Filler", layout="wide")
st.title("📄📊 Smart Placeholder Filler for DOCX & PPTX")

# 📘 How to Use
with st.expander("ℹ️ How to Use This App", expanded=True):
    st.markdown("""
    **Steps to generate a DOCX or PPTX document:**

    1. Upload your `.docx` or `.pptx` template with `{placeholders}`.
    2. Enter values for key fields like `{CUSTOMER_NAME}`, `{CITY NAME}`.
    3. Upload files or type values for all other placeholders:
       - `.docx`, `.txt`, `.pptx`: Extracts and inserts text
       - `.xlsx`: Inserts as Word table
       - `.jpg`, `.png`: Embedded as image (DOCX only)
    4. Click **🛠 Generate Document** to build and download.

    💡 **Tips:**
    - If image upload fails, type a manual description instead.
    - Keep placeholder names clean and consistent.
    """)

TEXT_ONLY_PLACEHOLDERS = {
    "CUSTOMER_NAME", "CITY NAME", "SA-NAME", "SA_EMAIL", "RAX_TEAM", "PARTNER_NAME"
}
today = date.today().strftime("%Y%m%d")

# Upload template
template_file = st.file_uploader("📁 Upload DOCX or PPTX template", type=["docx", "pptx"])
doc_type = st.selectbox("📄 Type of Document", ["Solution Proposal", "Migration Plan", "Report", "Presentation"])
customer_name = st.text_input("👤 Customer Name")

if template_file and customer_name:
    is_docx = template_file.name.endswith(".docx")
    is_pptx = template_file.name.endswith(".pptx")
    uploads = {}

    # Extract placeholders
    text_blocks = []
    if is_docx:
        doc = Document(template_file)
        text_blocks = [p.text for p in doc.paragraphs]
    elif is_pptx:
        prs = Presentation(template_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_blocks.append(shape.text)

    raw_placeholders = re.findall(r"\{[^}]+\}", "\n".join(text_blocks))
    placeholders = list(dict.fromkeys([f"{{{ph.strip('{}').strip()}}}" for ph in raw_placeholders]))

    # Step 1: Manual text input
    st.markdown("### ✏️ Enter Values for Key Fields")
    for ph in placeholders:
        base = ph.strip("{}").strip()
        if base in TEXT_ONLY_PLACEHOLDERS:
            val = st.text_input(f"✏️ {ph}", key=f"text_{base}")
            if val.strip():
                uploads[ph] = val.strip()

    # Step 2: Uploads or text for remaining placeholders
    st.markdown("### 📎 Upload Files or Enter Text for Other Placeholders")
    for ph in placeholders:
        base = ph.strip("{}").strip()
        if base not in TEXT_ONLY_PLACEHOLDERS:
            col1, col2 = st.columns(2)
            with col1:
                st.markdown("*Supported: .docx, .txt, .xlsx, .pptx, .jpg, .png*")
                file = st.file_uploader(
                    f"📎 Upload for {ph}",
                    type=["docx", "txt", "xlsx", "pptx", "jpg", "jpeg", "png"],
                    key=f"file_{base}"
                )
            with col2:
                text = st.text_area(f"✏️ Or enter value for {ph}", key=f"text_{base}")

            if file:
                ext = file.name.lower().split(".")[-1]
                if ext in ["jpg", "jpeg", "png"]:
                    img_bytes = BytesIO(file.read())
                    uploads[ph] = img_bytes
                elif ext == "xlsx":
                    df = pd.read_excel(file)
                    uploads[ph] = df
                elif ext == "docx":
                    d = Document(file)
                    uploads[ph] = "\n".join(p.text for p in d.paragraphs)
                elif ext == "pptx":
                    p = Presentation(file)
                    uploads[ph] = "\n".join(
                        shape.text for slide in p.slides for shape in slide.shapes if hasattr(shape, "text")
                    )
                elif ext == "txt":
                    uploads[ph] = file.read().decode("utf-8")
                else:
                    uploads[ph] = f"[Unsupported file type: {file.name}]"
            elif text.strip():
                uploads[ph] = text.strip()

    # Step 3: Generate output
    if st.button("🛠️ Generate Document"):
        final_filename = f"{customer_name}_{doc_type.replace(' ', '_')}_{today}"
        buffer = BytesIO()

        if is_docx:
            doc = Document(template_file)
            for para in doc.paragraphs:
                for ph, val in uploads.items():
                    if ph in para.text:
                        para.text = para.text.replace(ph, "")
                        if isinstance(val, BytesIO):
                            val.seek(0)
                            with NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                                tmp.write(val.read())
                                tmp.flush()
                                new_para = doc.add_paragraph()
                                new_para.add_run().add_picture(tmp.name, width=Inches(4))
                                os.unlink(tmp.name)
                        elif isinstance(val, pd.DataFrame):
                            table = doc.add_table(rows=1, cols=len(val.columns))
                            hdr_cells = table.rows[0].cells
                            for i, col in enumerate(val.columns):
                                hdr_cells[i].text = col
                            for _, row in val.iterrows():
                                row_cells = table.add_row().cells
                                for i, cell in enumerate(row):
                                    row_cells[i].text = str(cell)
                        else:
                            para.add_run(str(val))
            doc.save(buffer)
            st.success("✅ DOCX generated!")
            st.download_button(
                "📥 Download DOCX",
                buffer.getvalue(),
                file_name=final_filename + ".docx"
            )

        elif is_pptx:
            prs = Presentation(template_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        for ph, val in uploads.items():
                            if ph in shape.text:
                                shape.text = shape.text.replace(ph, str(val))
            prs.save(buffer)
            st.success("✅ PPTX generated!")
            st.download_button(
                "📥 Download PPTX",
                buffer.getvalue(),
                file_name=final_filename + ".pptx"
            )
